from pptx import Presentation

def remove_all_slides(prs):
    xml_slides = prs.slides._sldIdLst  
    slide_ids = list(xml_slides)
    for slide_id in slide_ids:
        xml_slides.remove(slide_id)

def create_presentation(template_file, slides_content):
    prs = Presentation(template_file)
    remove_all_slides(prs)

    for i, slide_data in enumerate(slides_content):
        if i == 0:
            layout = prs.slide_layouts[0]  # Title Slide
        else:
            layout = prs.slide_layouts[1]  # Title + Content
        
        slide = prs.slides.add_slide(layout)

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data["title"]

        # Body text (only if layout has a content placeholder)
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_data["content"]

    return prs

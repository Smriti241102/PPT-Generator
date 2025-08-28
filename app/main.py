import os
import json
import tempfile
from fastapi import FastAPI, File, UploadFile, Form, Request
from fastapi.responses import FileResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
from pptx import Presentation
from app.ppt_generator import create_presentation
from pptx.util import Pt
import httpx

app = FastAPI()
# Serve static frontend
app.mount("/static", StaticFiles(directory="app/static"), name="static")

INDEX_HTML = ""  # not used; static file will be served

# Helper: call OpenAI ChatCompletion (example). We treat provider generically but include OpenAI and Anthropic basic flows.
async def call_openai_chat(openai_key: str, provider:str, prompt: str, model: str = "gpt-4o-mini") -> str:
    if provider == "openai":
        url = "https://aipipe.org/openai/v1/chat/completions"
    elif provider == "gemini":
        url = "http://aipipe.org/openrouter/v1/chat/completions"
    else:
        raise ValueError("Unsupported provider")
    
    headers = {"Authorization": f"Bearer {openai_key}", "Content-Type": "application/json"}
    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content": "You convert user text into a JSON array of slides. Output only valid JSON."},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.2,
        "max_tokens": 1200
    }
    async with httpx.AsyncClient(timeout=120) as client:
        r = await client.post(url, headers=headers, json=payload)
        if r.status_code != 200:
            return f"Proxy error {r.status_code}: {r.text}"
        data = r.json()
    # The model's content
    return data["choices"][0]["message"]["content"]

# A simple LLM-to-slide JSON prompt builder
def build_prompt(text: str, guidance: str = "") -> str:
    system = (
        "Convert the following text (markdown or prose) into a JSON object named 'slides'."
        " The JSON object should look like: { \"slides\": [{\"title\":..., \"content\": [..], \"notes\": \"..\"}, ...] }"
        " Do NOT add any extra text before/after the JSON. Content array items should be short bullet lines (max 20 words)."
    )
    guidance_line = f"Use this guidance: {guidance}\n" if guidance else ""
    prompt = f"{system}\n{guidance_line}\nINPUT:\n{text}\n\nRespond with JSON only."
    return prompt

# Extract JSON robustly (try to find first { ... } block)
def extract_json(text: str):
    text = text.strip()
    # naive extraction: find first { and last }
    try:
        start = text.index("{")
        end = text.rindex("}")
        json_text = text[start:end+1]
        return json.loads(json_text)
    except Exception:
        # last resort: try direct json loads
        return json.loads(text)


@app.post("/generate")
async def generate_ppt(
    text: str = Form(...),
    guidance: str = Form("") ,
    provider: str = Form("openai"),
    model: str = Form("gpt-4o-mini"),
    api_key: str = Form(...),
    template: UploadFile = File(...)
):
    # Save uploaded template temporarily
    tmp_template = tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(template.filename)[1])
    contents = await template.read()
    tmp_template.write(contents)
    tmp_template.flush()

    # Build prompt and call provider
    prompt = build_prompt(text, guidance)
    # Only OpenAI implemented in this example. You can add Anthropic/Gemini by adding more helper functions.
    raw = await call_openai_chat(api_key, provider, prompt, model=model)
    

    # extract JSON safely
    slides_obj = extract_json(raw)
    if "slides" not in slides_obj:
        raise ValueError("LLM did not return a 'slides' key in the JSON. Response received: " + raw[:500])

    # Open the uploaded template presentation
    prs = Presentation(tmp_template.name)
    # prs_copy = Presentation(tmp_template.name)

    
    # Choose a layout: prefer title + content when available
    def choose_layout(prs):
        # Try to find layout with two placeholders (title + body)
        for i, l in enumerate(prs.slide_layouts):
            ph = [p for p in l.placeholders]
            if len(ph) >= 2:
                return i
        return 0


    base_layout_index = choose_layout(prs)
    slides_to_remove = len(prs.slides)

    def remove_first_n_slides(prs, n):
        for i in range(n):
            slide_id = prs.slides._sldIdLst[0]
            prs.slides._sldIdLst.remove(slide_id)

    print(len(slides_obj["slides"]))
    
    
    # Create slides
    for i,s in enumerate(slides_obj["slides"]):
        title = s.get("title", "")
        content_items = s.get("content", [])
        notes = s.get("notes", "")
        
        layout = prs.slide_layouts[base_layout_index]
        
        
        slide = prs.slides.add_slide(layout)
        # Title
        try:
            if slide.shapes.title:
                slide.shapes.title.text = title
        except Exception:
            print('exception occurred')
            pass

        # Body placeholder often at index 1
        body_placeholder = None
        for shp in slide.shapes:
            if shp.is_placeholder and shp.placeholder_format.idx != 0:
                body_placeholder = shp
                break
        if not body_placeholder:
            # fallback: add a textbox
            left = top = width = height = None
            from pptx.util import Inches
            txbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
            tf = txbox.text_frame
            for i, b in enumerate(content_items):
                if i == 0:
                    tf.text = b
                else:
                    p = tf.add_paragraph()
                    p.text = b
                    p.level = 0
        else:
            tf = body_placeholder.text_frame
            tf.clear()
            first = True
            for b in content_items:
                if first:
                    tf.text = b
                    first = False
                else:
                    p = tf.add_paragraph()
                    p.text = b
                    p.level = 0

        # Speaker notes
        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass

        print(len(prs.slides))

    print('final', len(prs.slides))
    remove_first_n_slides(prs, slides_to_remove)
    print('final2', len(prs.slides))
    # Save to temp file and return
    out_tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(out_tmp.name)
    return FileResponse(out_tmp.name, filename="generated_presentation.pptx", media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")


@app.get("/")
async def root():
    html_path = os.path.join("app", "static", "index.html")
    with open(html_path, "r", encoding="utf-8") as f:
        return HTMLResponse(f.read())
